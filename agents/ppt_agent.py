from __future__ import annotations

from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt

from utils.schemas import MonthlySynthesis, WeeklyReport


class PPTAgent:
    navy = RGBColor(18, 35, 56)
    teal = RGBColor(0, 122, 125)
    red = RGBColor(192, 63, 63)
    amber = RGBColor(210, 151, 46)
    green = RGBColor(38, 142, 92)
    grey = RGBColor(92, 103, 115)

    def _title(self, slide, headline: str, message: str = "") -> None:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(248, 250, 252)
        box = slide.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(12.1), Inches(0.55))
        p = box.text_frame.paragraphs[0]
        p.text = headline
        p.font.size = Pt(25)
        p.font.bold = True
        p.font.color.rgb = self.navy
        if message:
            sub = slide.shapes.add_textbox(Inches(0.57), Inches(0.92), Inches(12.0), Inches(0.45))
            sp = sub.text_frame.paragraphs[0]
            sp.text = message
            sp.font.size = Pt(12)
            sp.font.color.rgb = self.grey

    def _bullets(self, slide, title: str, items: list[str], x: float, y: float, w: float, h: float, size: int = 13) -> None:
        head = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.3))
        hp = head.text_frame.paragraphs[0]
        hp.text = title
        hp.font.bold = True
        hp.font.size = Pt(13)
        hp.font.color.rgb = self.teal
        box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.35), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        for idx, item in enumerate(items or ["No evidence available."]):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = item[:230]
            p.font.size = Pt(size)
            p.font.color.rgb = self.navy
            p.space_after = Pt(6)

    def _card(self, slide, label: str, value: str, x: float, y: float, color: RGBColor) -> None:
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(2.15), Inches(0.75))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = color
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(8)
        p2.font.color.rgb = RGBColor(255, 255, 255)

    def _top_risks(self, reports: list[WeeklyReport]) -> list[str]:
        risks = []
        for report in reports:
            for risk in report.features.risk_register[:3]:
                risks.append(
                    f"{risk['risk']} | Impact: {risk['business_impact']} | Likelihood: {risk['likelihood']} | Owner: {risk['owner']} | Mitigation: {risk['mitigation']}"
                )
        return risks[:5]

    def run(self, synthesis: MonthlySynthesis, reports: list[WeeklyReport], output_root: str | Path = "output") -> MonthlySynthesis:
        root = Path(output_root) / "ppt"
        root.mkdir(parents=True, exist_ok=True)
        reports = [r for r in reports if r.project.project_name != "Unknown Project"]
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        latest = sorted(reports, key=lambda r: r.generated_at)
        best = max(latest, key=lambda r: r.scoring.score, default=None)
        riskiest = min(latest, key=lambda r: r.scoring.score, default=None)
        red = sum(1 for r in latest if r.scoring.rag == "Red")
        amber = sum(1 for r in latest if r.scoring.rag == "Amber")
        green = sum(1 for r in latest if r.scoring.rag == "Green")

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._title(slide, "Executive Summary", "Overall Portfolio Health and recommended leadership action")
        self._card(slide, "Green", str(green), 0.7, 1.35, self.green)
        self._card(slide, "Amber", str(amber), 3.0, 1.35, self.amber)
        self._card(slide, "Red", str(red), 5.3, 1.35, self.red)
        self._bullets(slide, "Key Message", synthesis.portfolio_overview + synthesis.executive_recommendations[:2], 0.8, 2.45, 11.7, 3.8, 15)

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._title(slide, "Portfolio Health", "Health distribution with executive narrative")
        chart_data = CategoryChartData()
        chart_data.categories = ["Green", "Amber", "Red"]
        chart_data.add_series("RAG", (green, amber, red))
        slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(0.8), Inches(1.55), Inches(4.6), Inches(3.8), chart_data)
        self._bullets(slide, "Narrative", synthesis.portfolio_overview + synthesis.leadership_attention[:3], 6.0, 1.55, 6.4, 3.9, 14)

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._title(slide, "Trend Analysis", "What improved, what worsened, and which issues are recurring")
        self._bullets(slide, "Trend Signals", synthesis.portfolio_trends[:7], 0.8, 1.3, 5.8, 5.3, 13)
        self._bullets(slide, "Improving / Deteriorating", [f"Improving: {', '.join(synthesis.projects_improving) or 'None'}", f"Deteriorating: {', '.join(synthesis.projects_deteriorating or synthesis.projects_declining) or 'None'}"] + synthesis.sentiment_trends, 7.0, 1.3, 5.4, 5.3, 13)

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._title(slide, "Emerging Risks", "Top risks with impact, likelihood, mitigation and owner")
        self._bullets(slide, "Top Five Risks", self._top_risks(latest) or synthesis.emerging_risks, 0.75, 1.25, 12.0, 5.6, 11)

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._title(slide, "Leadership Decisions", "Decisions required to improve delivery confidence")
        decisions = [
            f"Decision: {rec} | Owner: Executive sponsor / PMO | Timeline: This week | Expected benefit: Removes delivery ambiguity."
            for rec in synthesis.executive_recommendations[:5]
        ]
        self._bullets(slide, "Decision Log", decisions, 0.75, 1.25, 12.0, 5.6, 12)

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._title(slide, "Project Spotlight", "Best performing and highest-risk projects")
        best_lines = [f"{best.project.project_name}: {best.scoring.rag} ({best.scoring.score}). {best.reasoning.executive_summary}" if best else "No best project available."]
        risk_lines = [f"{riskiest.project.project_name}: {riskiest.scoring.rag} ({riskiest.scoring.score}). {riskiest.reasoning.executive_summary}" if riskiest else "No high-risk project available."]
        self._bullets(slide, "Best Performing", best_lines, 0.8, 1.35, 5.7, 4.8, 12)
        self._bullets(slide, "Highest Risk", risk_lines, 7.0, 1.35, 5.5, 4.8, 12)

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._title(slide, "Appendix", "Supporting weekly metrics")
        appendix = [
            f"{snap['project_name']} | Week {snap['week_number']} | {snap['date']} | {snap['rag']} | Score {snap['health_score']} | Completion {snap['completion']}%"
            for snap in synthesis.project_snapshots
        ]
        self._bullets(slide, "Metrics", appendix, 0.75, 1.25, 12.0, 5.6, 11)

        path = root / f"{datetime.now():%Y%m%d_%H%M%S}_pulsepm_executive_report.pptx"
        prs.save(path)
        synthesis.ppt_path = str(path)
        return synthesis
